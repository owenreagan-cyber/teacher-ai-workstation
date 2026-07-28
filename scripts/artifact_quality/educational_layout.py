"""Grade 4 educational layout heuristics — conservative, WARN-first, never claims pedagogy."""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

import fitz

from .models import CheckStatus, PreflightReport
from .profiles import ArtifactProfile, EducationalLayoutSpec
from .visual_geometry import PageVisualMetrics


@dataclass
class TextSpanInfo:
    text: str
    size: float
    font: str
    flags: int
    y0: float
    y1: float


@dataclass
class PageStructure:
    page_number: int
    spans: list[TextSpanInfo]
    char_count: int
    question_count: int
    paragraph_lengths: list[int]
    heading_sizes: list[float]
    body_sizes: list[float]
    direction_sizes: list[float]
    has_drawing_separators: bool
    block_tops: list[float]
    blank_line_count: int


@dataclass
class EducationalLayoutReport:
    categories: dict[str, str] = field(default_factory=dict)
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = field(default_factory=list)
    educational_score: float = 100.0
    notes: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "categories": self.categories,
            "educational_score": round(self.educational_score, 1),
            "notes": self.notes,
            "finding_count": len(self.findings),
        }


def _normalize_font(name: str) -> str:
    base = name.split("+")[-1].split(",")[0].strip().lower()
    for prefix in ("times", "arial", "helv", "cour", "calibri", "georgia"):
        if prefix in base:
            return prefix
    return base or "unknown"


def _is_bold(flags: int) -> bool:
    return bool(flags & 2**4)


def _extract_page_structure(page: fitz.Page, page_number: int) -> PageStructure:
    spans: list[TextSpanInfo] = []
    block_tops: list[float] = []
    try:
        drawings = page.get_drawings()
        has_separators = any(
            any(isinstance(item, (list, tuple)) and item[0] == "l" for item in d.get("items", []))
            for d in drawings
        )
    except Exception:  # noqa: BLE001
        has_separators = False

    raw = page.get_text("dict")
    for block in raw.get("blocks", []):
        if block.get("type") != 0:
            continue
        block_tops.append(block.get("bbox", [0, 0, 0, 0])[1])
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = span.get("text", "").strip()
                if not text:
                    continue
                bbox = span.get("bbox", (0, 0, 0, 0))
                spans.append(
                    TextSpanInfo(
                        text=text,
                        size=float(span.get("size", 0)),
                        font=_normalize_font(str(span.get("font", ""))),
                        flags=int(span.get("flags", 0)),
                        y0=bbox[1],
                        y1=bbox[3],
                    )
                )

    full_text = page.get_text("text")
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", full_text) if p.strip()]
    paragraph_lengths = [len(re.sub(r"\s+", "", p)) for p in paragraphs] or [len(re.sub(r"\s+", "", full_text))]
    questions = len(re.findall(r"(?m)^\s*\d+[\.)]\s", full_text))
    char_count = len(re.sub(r"\s+", "", full_text))

    sizes = [s.size for s in spans if s.size > 0]
    heading_sizes = [s.size for s in spans if s.size >= 16]
    body_sizes = [s.size for s in spans if 0 < s.size < 16]
    direction_sizes = [s.size for s in spans if "direction" in s.text.lower()]

    blank_lines = len(re.findall(r"(?m)^\s*$", full_text))

    return PageStructure(
        page_number=page_number,
        spans=spans,
        char_count=char_count,
        question_count=questions,
        paragraph_lengths=paragraph_lengths,
        heading_sizes=heading_sizes,
        body_sizes=body_sizes,
        direction_sizes=direction_sizes,
        has_drawing_separators=has_separators,
        block_tops=sorted(block_tops),
        blank_line_count=blank_lines,
    )


def _extract_pages(doc: fitz.Document) -> list[PageStructure]:
    return [_extract_page_structure(page, i + 1) for i, page in enumerate(doc)]


def _category_status(findings: Iterable[tuple[CheckStatus, str, str | None, int | None]]) -> str:
    statuses = [f[0] for f in findings]
    if any(s == CheckStatus.FAIL for s in statuses):
        return "FAIL"
    if any(s == CheckStatus.WARN for s in statuses):
        return "WARN"
    return "PASS"


def _add(
    report: EducationalLayoutReport,
    status: CheckStatus,
    message: str,
    *,
    details: str | None = None,
    page: int | None = None,
) -> None:
    report.findings.append((status, message, details, page))


def analyze_typography(
    pages: list[PageStructure],
    profile: ArtifactProfile,
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    min_body = max(profile.requirements.min_body_font_pt, spec.min_body_font_pt)
    min_direction = spec.min_direction_font_pt
    fail_floor = spec.unreadable_body_font_pt

    all_body: list[float] = []
    all_fonts: set[str] = set()
    heading_by_page: dict[int, list[float]] = {}

    for pg in pages:
        all_body.extend(pg.body_sizes)
        for span in pg.spans:
            all_fonts.add(span.font)
        if pg.heading_sizes:
            heading_by_page[pg.page_number] = pg.heading_sizes

        small_body = [size for size in pg.body_sizes if size < min_body]
        unreadable = [size for size in pg.body_sizes if size < fail_floor]
        for size in unreadable[:1]:
            findings.append(
                (
                    CheckStatus.FAIL,
                    f"Page {pg.page_number} body text may be objectively unreadable ({size:.1f} pt)",
                    f"Below unreadable threshold ({fail_floor:.1f} pt).",
                    pg.page_number,
                )
            )
        if unreadable:
            continue
        if small_body:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} body text unusually small (min {min(small_body):.1f} pt)",
                    f"Grade 4 profile expects at least {min_body:.1f} pt.",
                    pg.page_number,
                )
            )

        for size in pg.direction_sizes:
            if size < min_direction:
                findings.append(
                    (
                        CheckStatus.WARN,
                        f"Page {pg.page_number} directions may blend with body text ({size:.1f} pt)",
                        f"Consider at least {min_direction:.1f} pt for directions.",
                        pg.page_number,
                    )
                )
                break

    if len(all_fonts) > spec.max_font_families:
        findings.append(
            (
                CheckStatus.WARN,
                f"Multiple unrelated font families detected ({len(all_fonts)})",
                f"Families: {', '.join(sorted(all_fonts))}",
                None,
            )
        )

    heading_sizes_flat = [s for sizes in heading_by_page.values() for s in sizes]
    if heading_sizes_flat and max(heading_sizes_flat) - min(heading_sizes_flat) > spec.max_heading_size_spread_pt:
        findings.append(
            (
                CheckStatus.WARN,
                "Heading sizes appear inconsistent across pages",
                f"Spread {max(heading_sizes_flat) - min(heading_sizes_flat):.1f} pt exceeds {spec.max_heading_size_spread_pt:.1f} pt.",
                None,
            )
        )

    return findings


def analyze_cognitive_load(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        if pg.char_count > spec.max_chars_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} text density is high ({pg.char_count} chars)",
                    "Large uninterrupted text may overwhelm Grade 4 readers.",
                    pg.page_number,
                )
            )
        long_paras = [n for n in pg.paragraph_lengths if n > spec.max_paragraph_chars]
        longest_span_run = max((len(s.text) for s in pg.spans), default=0)
        if long_paras or longest_span_run > spec.max_paragraph_chars:
            longest = max(long_paras) if long_paras else longest_span_run
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} has large paragraph block(s)",
                    f"Longest block ~{longest} chars; consider chunking.",
                    pg.page_number,
                )
            )
        if pg.question_count > spec.max_questions_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} has many simultaneous tasks ({pg.question_count} numbered items)",
                    f"Profile suggests no more than {spec.max_questions_per_page} per page.",
                    pg.page_number,
                )
            )
    return findings


def analyze_chunking(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        has_heading = bool(pg.heading_sizes)
        has_separator = pg.has_drawing_separators or pg.blank_line_count >= 2
        gaps = [
            pg.block_tops[i + 1] - pg.block_tops[i]
            for i in range(len(pg.block_tops) - 1)
            if pg.block_tops[i + 1] - pg.block_tops[i] > 8
        ]
        weak = pg.char_count > spec.min_chars_for_chunking_check and not has_heading and not has_separator and len(gaps) < 2
        if weak:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} chunking appears weak",
                    "Continuous content without clear heading, divider, or spacing.",
                    pg.page_number,
                )
            )
    return findings


def analyze_writing_space_educational(
    page_metrics: list[PageVisualMetrics],
    profile: ArtifactProfile,
    subject: str | None,
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    ws_min, ws_max = profile.visual_geometry.writing_space_range
    subj = (subject or "").lower().replace("_", "-")

    for m in page_metrics:
        if m.writing_space_percent <= 0:
            continue
        expected = spec.subject_writing_space.get(subj, (ws_min, ws_max))
        if m.writing_space_percent < expected[0]:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {m.page_number} may lack sufficient writing space ({m.writing_space_percent:.0f}%)",
                    f"Expected roughly {expected[0]:.0f}–{expected[1]:.0f}% for {subj or profile.name}.",
                    m.page_number,
                )
            )
        elif m.writing_space_percent > expected[1] + 15:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {m.page_number} writing-space estimate unusually high ({m.writing_space_percent:.0f}%)",
                    "Confirm intentional workspace vs decorative empty area.",
                    m.page_number,
                )
            )
    return findings


def analyze_directions(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        full = " ".join(s.text for s in pg.spans)
        match = re.search(r"directions?[:\s]+(.{20,800})", full, re.IGNORECASE)
        if not match:
            continue
        block = match.group(1)
        words = len(block.split())
        sentences = max(len(re.findall(r"[.!?]", block)), 1)
        if words > spec.max_direction_words:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} directions may be overly long ({words} words)",
                    "Consider bullets or numbered steps with action verbs.",
                    pg.page_number,
                )
            )
        if sentences == 1 and words > spec.max_direction_single_paragraph_words:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} multi-step instructions appear in one paragraph",
                    "Split dense instruction blocks for Grade 4 clarity.",
                    pg.page_number,
                )
            )
    return findings


def analyze_worksheet_layout(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    activity_markers = 0
    for pg in pages:
        formats = 0
        if pg.question_count:
            formats += 1
        if pg.has_drawing_separators:
            formats += 1
        if any("show work" in s.text.lower() or "answer" in s.text.lower() for s in pg.spans):
            formats += 1
        if formats > spec.max_activity_formats_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} mixes many activity formats",
                    "Consider grouping similar question types.",
                    pg.page_number,
                )
            )
        activity_markers += formats
    return findings


def analyze_guided_notes_layout(
    pages: list[PageStructure],
    page_metrics: list[PageVisualMetrics],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        blankish = sum(1 for s in pg.spans if re.fullmatch(r"[_\.\-—]+", s.text))
        text_chars = sum(len(s.text) for s in pg.spans)
        if text_chars > 0 and blankish / max(len(pg.spans), 1) < spec.min_fill_in_ratio:
            long_copy = [s for s in pg.spans if len(s.text) > 60 and s.size < 14]
            if long_copy:
                findings.append(
                    (
                        CheckStatus.WARN,
                        f"Page {pg.page_number} guided notes may require excessive copying",
                        "Large text blocks with few fill-in blanks detected.",
                        pg.page_number,
                    )
                )
    for m in page_metrics:
        if m.writing_space_percent > 50 and m.text_coverage_percent > 25:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {m.page_number} may require large blank transcription areas",
                    "Confirm students are not copying full paragraphs.",
                    m.page_number,
                )
            )
    return findings


def analyze_assessment_layout(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    q_styles: set[str] = set()
    for pg in pages:
        if pg.question_count > spec.max_assessment_questions_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} assessment density is high ({pg.question_count} items)",
                    "Consider splitting across pages or front/back.",
                    pg.page_number,
                )
            )
        if re.search(r"\?\s*$", " ".join(s.text for s in pg.spans), re.MULTILINE):
            q_styles.add("question")
        if re.search(r"=\s*_{3,}", " ".join(s.text for s in pg.spans)):
            q_styles.add("blank")
        if re.search(r"\bcircle\b|\bselect\b", " ".join(s.text for s in pg.spans), re.I):
            q_styles.add("choice")
    if len(q_styles) > spec.max_question_style_variety:
        findings.append(
            (
                CheckStatus.WARN,
                "Assessment mixes many question styles on one artifact",
                f"Styles detected: {', '.join(sorted(q_styles))}",
                None,
            )
        )
    return findings


def analyze_reading_layout(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        lines = [s for s in pg.spans if len(s.text) > 40]
        if len(lines) > spec.max_reading_lines_without_break:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} passage may be visually crowded",
                    "Long lines without clear question separation.",
                    pg.page_number,
                )
            )
        if pg.question_count and pg.char_count > spec.max_reading_passage_chars:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} questions may be embedded in dense passage text",
                    "Separate passage block from questions when possible.",
                    pg.page_number,
                )
            )
    return findings


def analyze_math_layout(
    pages: list[PageStructure],
    page_metrics: list[PageVisualMetrics],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        math_lines = sum(1 for s in pg.spans if re.search(r"\d+\s*[+\-×x*/=]", s.text))
        if math_lines > spec.max_math_problems_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} math problems appear crowded ({math_lines} expressions)",
                    "Allow vertical workspace between problems.",
                    pg.page_number,
                )
            )
    for m in page_metrics:
        if m.writing_space_percent < spec.min_math_workspace_percent and m.text_coverage_percent > 10:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {m.page_number} may lack computation workspace",
                    f"Writing-space estimate {m.writing_space_percent:.0f}% below math expectation.",
                    m.page_number,
                )
            )
    return findings


def analyze_shurley_layout(
    pages: list[PageStructure],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        for span in pg.spans:
            if len(span.text) > spec.max_shurley_sentence_chars and span.text.count(" ") > 8:
                findings.append(
                    (
                        CheckStatus.WARN,
                        f"Page {pg.page_number} Shurley sentence may wrap ({len(span.text)} chars)",
                        "Reserve vertical room for classification marks.",
                        pg.page_number,
                    )
                )
                break
        gaps = [pg.block_tops[i + 1] - pg.block_tops[i] for i in range(len(pg.block_tops) - 1)]
        if gaps and min(gaps) < spec.min_shurley_sentence_spacing_pt:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} sentence spacing may be tight for grammar marking",
                    "Classification marks may collide.",
                    pg.page_number,
                )
            )
    return findings


def analyze_diagram_layout(
    pages: list[PageStructure],
    page_metrics: list[PageVisualMetrics],
    spec: EducationalLayoutSpec,
) -> list[tuple[CheckStatus, str, str | None, int | None]]:
    findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    for pg in pages:
        tiny = [s for s in pg.spans if s.size < spec.min_diagram_label_font_pt and len(s.text) <= 20]
        if len(tiny) >= spec.max_tiny_labels_per_page:
            findings.append(
                (
                    CheckStatus.WARN,
                    f"Page {pg.page_number} diagram labels may be too small",
                    f"{len(tiny)} labels below {spec.min_diagram_label_font_pt:.1f} pt.",
                    pg.page_number,
                )
            )
    for m in page_metrics:
        if m.drawing_coverage_percent > 20 and m.text_coverage_percent < 3 and m.visible_ink_percent > 15:
            if m.visible_ink_percent > spec.max_diagram_crowding_ink_percent:
                findings.append(
                    (
                        CheckStatus.WARN,
                        f"Page {m.page_number} diagram region appears crowded",
                        "Confirm labels, legends, and timelines remain readable.",
                        m.page_number,
                    )
                )
    return findings


def analyze_pptx_educational(path: Path, profile: ArtifactProfile, spec: EducationalLayoutSpec) -> EducationalLayoutReport:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    report = EducationalLayoutReport()
    pres = Presentation(str(path))
    slide_findings: list[tuple[CheckStatus, str, str | None, int | None]] = []
    min_title = profile.slides.min_title_font_pt if profile.slides else spec.min_slide_title_font_pt
    min_body = profile.slides.min_body_font_pt if profile.slides else spec.min_slide_body_font_pt

    for idx, slide in enumerate(pres.slides, start=1):
        texts: list[str] = []
        sizes: list[float] = []
        bullets = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.level == 0 and para.text.strip():
                    bullets += 1
                for run in para.runs:
                    if run.text.strip():
                        texts.append(run.text.strip())
                        if run.font.size:
                            sizes.append(run.font.size.pt)
        body = " ".join(texts)
        if len(body.split()) > spec.max_slide_words:
            slide_findings.append(
                (
                    CheckStatus.WARN,
                    f"Slide {idx} reading load is high ({len(body.split())} words)",
                    "Projector slides should stay concise for Grade 4.",
                    idx,
                )
            )
        if bullets > spec.max_bullets_per_slide:
            slide_findings.append(
                (
                    CheckStatus.WARN,
                    f"Slide {idx} has excessive bullet count ({bullets})",
                    f"Profile suggests no more than {spec.max_bullets_per_slide}.",
                    idx,
                )
            )
        if sizes:
            if min(sizes) < min_body:
                slide_findings.append(
                    (
                        CheckStatus.WARN,
                        f"Slide {idx} projected text may be too small ({min(sizes):.1f} pt)",
                        f"Minimum body target {min_body:.1f} pt.",
                        idx,
                    )
                )
            title_sizes = [s for s in sizes if s >= min_title]
            if texts and not title_sizes:
                slide_findings.append(
                    (
                        CheckStatus.WARN,
                        f"Slide {idx} title prominence may be weak",
                        f"Expected title at least {min_title:.1f} pt.",
                        idx,
                    )
                )
        concept_markers = len(re.findall(r"\b(and|also|another|first|second|third)\b", body, re.I))
        if concept_markers > spec.max_slide_concept_markers:
            slide_findings.append(
                (
                    CheckStatus.WARN,
                    f"Slide {idx} may contain multiple unrelated concepts",
                    "One-core-idea-per-slide heuristic flagged multiple topic markers.",
                    idx,
                )
            )

    report.findings.extend(slide_findings)
    report.categories = {
        "Typography": "PASS",
        "Visual Chunking": "PASS",
        "Writing Space": "PASS",
        "Text Density": _category_status(slide_findings),
        "Directions": "PASS",
        "Grade 4 Readability": _category_status(slide_findings),
        "Presentation Visibility": _category_status(slide_findings),
        "Manual Review": "REQUIRED",
    }
    report.educational_score = _compute_educational_score(report.findings, fail_present=False)
    return report


def analyze_pdf_educational_layout(
    doc: fitz.Document,
    profile: ArtifactProfile,
    subject: str | None,
    page_metrics: list[PageVisualMetrics],
) -> EducationalLayoutReport:
    spec = profile.educational_layout
    pages = _extract_pages(doc)
    layout = EducationalLayoutReport()

    buckets: dict[str, list[tuple[CheckStatus, str, str | None, int | None]]] = {
        "Typography": analyze_typography(pages, profile, spec),
        "Text Density": analyze_cognitive_load(pages, spec),
        "Visual Chunking": analyze_chunking(pages, spec),
        "Writing Space": analyze_writing_space_educational(page_metrics, profile, subject, spec),
        "Directions": analyze_directions(pages, spec),
        "Grade 4 Readability": [],
    }

    name = profile.name.lower()
    subj = (subject or "").lower().replace("_", "-")

    if "worksheet" in name:
        buckets["Grade 4 Readability"].extend(analyze_worksheet_layout(pages, spec))
    if "guided-notes" in name:
        buckets["Grade 4 Readability"].extend(analyze_guided_notes_layout(pages, page_metrics, spec))
    if "quiz" in name:
        buckets["Grade 4 Readability"].extend(analyze_assessment_layout(pages, spec))
    if subj == "reading":
        buckets["Grade 4 Readability"].extend(analyze_reading_layout(pages, spec))
    if subj == "math":
        buckets["Writing Space"].extend(analyze_math_layout(pages, page_metrics, spec))
    if subj == "shurley":
        buckets["Grade 4 Readability"].extend(analyze_shurley_layout(pages, spec))
    if subj in {"history", "science"}:
        buckets["Grade 4 Readability"].extend(analyze_diagram_layout(pages, page_metrics, spec))

    for items in buckets.values():
        layout.findings.extend(items)

    layout.categories = {k: _category_status(v) if v else "PASS" for k, v in buckets.items()}
    layout.categories["Presentation Visibility"] = "PASS"
    layout.categories["Manual Review"] = "REQUIRED"
    fail_present = any(f[0] == CheckStatus.FAIL for f in layout.findings)
    layout.educational_score = _compute_educational_score(layout.findings, fail_present=fail_present)
    layout.notes.append("Educational heuristics identify probable layout issues; instructional approval remains manual.")
    return layout


def _compute_educational_score(
    findings: list[tuple[CheckStatus, str, str | None, int | None]],
    *,
    fail_present: bool,
) -> float:
    if fail_present:
        return 0.0
    warns = sum(1 for f in findings if f[0] == CheckStatus.WARN)
    return max(0.0, 100.0 - warns * 8.0)


def format_instructional_layout_section(layout: EducationalLayoutReport) -> str:
    lines = ["Instructional Layout", "------------------", ""]
    order = (
        "Typography",
        "Visual Chunking",
        "Writing Space",
        "Text Density",
        "Directions",
        "Grade 4 Readability",
        "Presentation Visibility",
        "Manual Review",
    )
    for key in order:
        status = layout.categories.get(key, "PASS")
        lines.append(f"{key}: {status}")
    lines.append("")
    return "\n".join(lines)


def apply_educational_layout(preflight: PreflightReport, layout: EducationalLayoutReport) -> None:
    preflight.educational_layout = layout.to_dict()
    preflight.add(CheckStatus.PASS, "Instructional layout analysis complete")
    for status, message, details, page in layout.findings:
        preflight.add(status, message, details=details, page=page)
    summary = format_instructional_layout_section(layout)
    preflight.add(CheckStatus.PASS, "Instructional layout summary", details=summary)


def compute_combined_quality_score(
    mechanical_pass_ratio: float,
    page_metrics: list[PageVisualMetrics],
    fail_count: int,
    warn_count: int,
    educational_score: float,
) -> dict[str, Any]:
    from .visual_geometry import compute_quality_score

    base = compute_quality_score(mechanical_pass_ratio, page_metrics, fail_count, warn_count)
    payload = base.to_dict()
    payload["educational_layout_score"] = round(educational_score, 1)
    payload["instructional_approval"] = "Manual Review Required"
    return payload
