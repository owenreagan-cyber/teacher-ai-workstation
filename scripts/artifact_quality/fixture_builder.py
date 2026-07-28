from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt


LETTER = (612, 792)
A4 = (595, 842)
# Safe inset for standard worksheet profile: margin + 0.35" boundary
SAFE_X = 72
SAFE_Y = 100
SAFE_BOTTOM_Y = 680


def _new_letter_doc() -> fitz.Document:
    doc = fitz.open()
    doc.new_page(width=LETTER[0], height=LETTER[1])
    return doc


def build_passing_worksheet(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Grade 4 Math Worksheet", fontsize=20)
    # Table grid with writing lines (inset within safe bounds)
    table_left = SAFE_X + 16
    table_right = 520
    table_top = SAFE_Y + 50
    for row in range(5):
        y = table_top + row * 36
        page.draw_line((table_left, y), (table_right, y), width=0.5)
    for col in range(3):
        x = table_left + col * ((table_right - table_left) // 3)
        page.draw_line((x, table_top), (x, table_top + 4 * 36), width=0.5)
    page.draw_line((table_right, table_top), (table_right, table_top + 4 * 36), width=0.5)
    page.draw_line((table_left, table_top + 4 * 36), (table_right, table_top + 4 * 36), width=0.5)
    col_width = (table_right - table_left) // 3
    for index in range(1, 9):
        col = (index - 1) % 3
        row = (index - 1) // 3
        x = table_left + col * col_width + 8
        y = table_top + row * 36 + 22
        page.insert_text((x, y), f"{index}.", fontsize=13)
    # Answer boxes
    for i in range(3):
        box_y = table_top + 4 * 36 + 20 + i * 50
        page.draw_rect(fitz.Rect(table_left, box_y, table_right, box_y + 36), width=0.75)
        page.insert_text((table_left + 8, box_y + 10), f"Show work {i + 1}:", fontsize=12.5)
        for line_y in range(box_y + 18, box_y + 34, 8):
            page.draw_line((table_left + 4, line_y), (table_right - 4, line_y), width=0.3)
    directions = (
        "Directions: Complete each problem. Show your work in the space provided. "
        "Check your answers carefully before turning in your paper."
    )
    page.insert_textbox(
        fitz.Rect(SAFE_X, SAFE_BOTTOM_Y - 60, 540, SAFE_BOTTOM_Y - 10),
        directions,
        fontsize=12.5,
    )
    page.insert_text((500, SAFE_BOTTOM_Y), "1", fontsize=12.5)
    doc.save(path)
    doc.close()


def build_passing_guided_notes(path: Path) -> None:
    doc = fitz.open()
    for page_no in (1, 2):
        page = doc.new_page(width=LETTER[0], height=LETTER[1])
        page.insert_text((SAFE_X + 8, SAFE_Y), f"Guided Notes — Page {page_no}", fontsize=18)
        page.insert_text((SAFE_X + 8, SAFE_Y + 40), "Main idea: sample summary line for testing.", fontsize=13)
        page.insert_text((SAFE_X + 8, SAFE_Y + 80), "Notes:", fontsize=13)
        for y in range(SAFE_Y + 120, SAFE_BOTTOM_Y - 120, 28):
            page.draw_line((SAFE_X + 8, y), (520, y), width=0.4)
        # Structured response box
        box_y = SAFE_BOTTOM_Y - 110
        page.draw_rect(fitz.Rect(SAFE_X + 8, box_y, 520, box_y + 80), width=0.75)
        page.insert_text((SAFE_X + 16, box_y + 8), "Response:", fontsize=12)
        page.insert_text((500, SAFE_BOTTOM_Y), str(page_no), fontsize=10)
    doc.save(path)
    doc.close()


def build_passing_teacher_key(path: Path) -> None:
    doc = fitz.open()
    for page_no in (1, 2):
        page = doc.new_page(width=LETTER[0], height=LETTER[1])
        page.insert_text((SAFE_X + 8, SAFE_Y), f"Teacher Key — Page {page_no}", fontsize=18)
        page.insert_text((SAFE_X + 8, SAFE_Y + 40), "Main idea: sample summary line for testing.", fontsize=13)
        page.insert_text((SAFE_X + 8, SAFE_Y + 80), "Answer: sample answer text.", fontsize=13)
        page.insert_text((500, SAFE_BOTTOM_Y), str(page_no), fontsize=10)
    doc.save(path)
    doc.close()


def build_passing_diagram_page(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Science Diagram", fontsize=18)
    cx, cy = 300, 380
    page.draw_circle(fitz.Point(cx, cy), 80, width=1.5)
    page.draw_line(fitz.Point(cx - 80, cy), fitz.Point(cx + 80, cy), width=0.75)
    page.draw_line(fitz.Point(cx, cy - 80), fitz.Point(cx, cy + 80), width=0.75)
    page.draw_line(fitz.Point(cx - 50, cy - 50), fitz.Point(cx + 50, cy + 50), width=0.5)
    page.draw_rect(fitz.Rect(SAFE_X, 520, 540, 620), width=0.75)
    page.insert_text((SAFE_X + 8, 530), "Label the diagram:", fontsize=12)
    doc.save(path)
    doc.close()


def build_warn_dense_worksheet(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Dense Worksheet", fontsize=16)
    shade = fitz.Rect(SAFE_X, SAFE_Y + 30, 520, SAFE_BOTTOM_Y - 40)
    page.draw_rect(shade, color=(0.2, 0.2, 0.2), fill=(0.92, 0.92, 0.92))
    y = SAFE_Y + 40
    while y < SAFE_BOTTOM_Y - 50:
        page.insert_text((SAFE_X + 8, y), "Dense sample instructional line for testing.", fontsize=9)
        page.draw_line((SAFE_X + 8, y + 12), (510, y + 12), width=0.4)
        y += 16
    doc.save(path)
    doc.close()


def build_warn_layout_shift_key(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=LETTER[0], height=LETTER[1])
    page.insert_text((SAFE_X + 8, SAFE_Y), "Teacher Key — Page 1", fontsize=18)
    page.insert_text((SAFE_X + 8, SAFE_Y + 40), "Main idea: sample summary line for testing.", fontsize=13)
    page.insert_text((SAFE_X + 8, SAFE_Y + 80), "Answer: sample answer text with extra teacher content.", fontsize=13)
    page.insert_text((SAFE_X + 8, SAFE_Y + 120), "Additional notes for teacher review only.", fontsize=13)
    page.insert_text((500, SAFE_BOTTOM_Y), "1", fontsize=10)
    page = doc.new_page(width=LETTER[0], height=LETTER[1])
    page.insert_text((SAFE_X + 8, SAFE_Y), "Teacher Key — Page 2", fontsize=18)
    page.insert_text((500, SAFE_BOTTOM_Y), "2", fontsize=10)
    doc.save(path)
    doc.close()


def build_passing_multipage(path: Path) -> None:
    doc = fitz.open()
    for page_no in range(1, 6):
        page = doc.new_page(width=LETTER[0], height=LETTER[1])
        page.insert_text((SAFE_X, SAFE_Y), f"Multi-page test — Page {page_no}", fontsize=16)
        for y in range(SAFE_Y + 60, SAFE_BOTTOM_Y - 40, 30):
            page.draw_line((SAFE_X, y), (540, y), width=0.3)
        page.insert_text((500, SAFE_BOTTOM_Y), str(page_no), fontsize=10)
    doc.save(path)
    doc.close()


def build_warn_low_utilization(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((80, 80), "Short title only", fontsize=14)
    doc.save(path)
    doc.close()


def build_warn_writing_space(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Minimal writing space worksheet", fontsize=16)
    page.insert_text((SAFE_X, SAFE_Y + 40), "1. Single question", fontsize=13)
    doc.save(path)
    doc.close()


def build_warn_bottom_gap(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((80, 80), "Section 1", fontsize=16)
    page.insert_text((80, 120), "A small amount of content near the top.", fontsize=13)
    doc.save(path)
    doc.close()


def build_warn_near_boundary(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((42, 42), "Edge content", fontsize=12)
    doc.save(path)
    doc.close()


def build_fail_a4(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=A4[0], height=A4[1])
    page.insert_text((60, 60), "A4 page sample", fontsize=14)
    doc.save(path)
    doc.close()


def build_fail_unsafe_edge(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((8, 400), "CLIPPED LEFT", fontsize=14)
    doc.save(path)
    doc.close()


def build_fail_blank_final_page(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=LETTER[0], height=LETTER[1])
    page.insert_text((SAFE_X, SAFE_Y), "Page 1 content", fontsize=14)
    doc.new_page(width=LETTER[0], height=LETTER[1])
    doc.save(path)
    doc.close()


def build_fail_mixed_sizes(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=LETTER[0], height=LETTER[1])
    page.insert_text((80, 80), "Letter page", fontsize=14)
    page = doc.new_page(width=A4[0], height=A4[1])
    page.insert_text((80, 80), "A4 page", fontsize=14)
    doc.save(path)
    doc.close()


def build_fail_placeholder(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((80, 80), "TODO finish this worksheet", fontsize=14)
    doc.save(path)
    doc.close()


def build_fail_shurley_wrap(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    sentence = (
        "The quick brown fox jumps over the lazy dog while the students classify every single word "
        "in this intentionally long sample sentence for testing."
    )
    page.insert_text((80, 120), sentence, fontsize=13)
    doc.save(path)
    doc.close()


def build_passing_html(path: Path) -> None:
    path.write_text(
        """<!DOCTYPE html>
<html>
<head>
<style>
@page { size: letter; margin: 0.5in; }
@media print {
  * { box-sizing: border-box; }
  .content { break-inside: avoid; page-break-inside: avoid; }
}
</style>
</head>
<body><div class="content"><h1>Reading Check</h1><p>Sample passage text.</p></div></body>
</html>
""",
        encoding="utf-8",
    )


def build_fail_html_no_print(path: Path) -> None:
    path.write_text("<html><body><p>No print styles</p></body></html>", encoding="utf-8")


def build_fail_html_overflow(path: Path) -> None:
    path.write_text(
        """<!DOCTYPE html>
<html><head><style>
.content { overflow: hidden; height: 200px; }
</style></head><body><div class="content worksheet">Academic content</div></body></html>
""",
        encoding="utf-8",
    )


def build_passing_docx(path: Path) -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.5)
    document.add_heading("Guided Notes DOCX", level=1)
    document.add_paragraph("Sample paragraph for structural validation.")
    document.save(path)


def build_fail_docx_a4(path: Path) -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.27)
    section.page_height = Inches(11.69)
    document.add_paragraph("A4 section")
    document.save(path)


def build_passing_pptx(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PInches(13.333)
    presentation.slide_height = PInches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(PInches(0.5), PInches(0.5), PInches(12), PInches(1))
    textbox.text_frame.text = "Projector Slide Title"
    body = slide.shapes.add_textbox(PInches(0.5), PInches(1.5), PInches(12), PInches(5))
    body.text_frame.text = "One core idea for classroom display."
    presentation.save(path)


def build_fail_pptx_out_of_bounds(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PInches(13.333)
    presentation.slide_height = PInches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    shape = slide.shapes.add_textbox(PInches(12.5), PInches(6.8), PInches(2), PInches(1))
    shape.text_frame.text = "Outside bounds"
    presentation.save(path)


def build_edu_warn_tiny_font(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Tiny Font Worksheet", fontsize=16)
    y = SAFE_Y + 40
    for i in range(12):
        page.insert_text((SAFE_X, y), f"{i + 1}. Problem with small body text.", fontsize=9.5)
        y += 18
    doc.save(path)
    doc.close()


def build_edu_warn_huge_paragraph(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Reading Passage", fontsize=18)
    paragraph = (
        "Directions: Read the entire passage carefully before answering. "
        + "This intentionally long single paragraph simulates a wall of text without visual breaks "
        * 8
        + "for Grade 4 readability testing."
    )
    page.insert_textbox(fitz.Rect(SAFE_X, SAFE_Y + 40, 520, 560), paragraph, fontsize=13)
    doc.save(path)
    doc.close()


def build_edu_warn_dense_slide(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PInches(13.333)
    presentation.slide_height = PInches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12), PInches(0.8))
    title.text_frame.text = "Dense Slide"
    body = slide.shapes.add_textbox(PInches(0.5), PInches(1.2), PInches(12), PInches(5.8))
    tf = body.text_frame
    bullets = [
        "First concept about ecosystems and also food chains",
        "Second concept about producers consumers decomposers",
        "Third concept about habitats and also adaptations",
        "Fourth concept about energy flow and also matter",
        "Fifth concept about human impact and also conservation",
        "Sixth concept about water cycle and also weather",
        "Seventh concept about soil layers and also erosion",
        "Eighth concept about plant parts and also photosynthesis",
    ]
    tf.text = bullets[0]
    for item in bullets[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = PPt(14)
    presentation.save(path)


def build_edu_warn_crowded_shurley(path: Path) -> None:
    doc = _new_letter_doc()
    page = doc[0]
    page.insert_text((SAFE_X, SAFE_Y), "Shurley Practice", fontsize=18)
    page.insert_text((SAFE_X, SAFE_Y + 36), "Classify each word in the sentence below.", fontsize=12)
    page.insert_text(
        (SAFE_X, SAFE_Y + 70),
        "The fourth grade students classified the vocabulary words.",
        fontsize=13,
    )
    page.insert_text(
        (SAFE_X, SAFE_Y + 100),
        "Another sample sentence with tight vertical spacing for testing.",
        fontsize=13,
    )
    doc.save(path)
    doc.close()


def ensure_all_fixtures(base: Path) -> None:
    builders = {
        base / "passing" / "worksheet-letter.pdf": build_passing_worksheet,
        base / "passing" / "guided-notes-two-page.pdf": build_passing_guided_notes,
        base / "passing" / "teacher-key-two-page.pdf": build_passing_teacher_key,
        base / "passing" / "diagram-minimal-text.pdf": build_passing_diagram_page,
        base / "passing" / "multipage-five.pdf": build_passing_multipage,
        base / "passing" / "printable.html": build_passing_html,
        base / "passing" / "guided-notes.docx": build_passing_docx,
        base / "passing" / "projector.pptx": build_passing_pptx,
        base / "warning" / "low-utilization.pdf": build_warn_low_utilization,
        base / "warning" / "bottom-gap.pdf": build_warn_bottom_gap,
        base / "warning" / "near-boundary.pdf": build_warn_near_boundary,
        base / "warning" / "dense-worksheet.pdf": build_warn_dense_worksheet,
        base / "warning" / "writing-space-low.pdf": build_warn_writing_space,
        base / "warning" / "layout-shift-key.pdf": build_warn_layout_shift_key,
        base / "warning" / "edu-tiny-font.pdf": build_edu_warn_tiny_font,
        base / "warning" / "edu-huge-paragraph.pdf": build_edu_warn_huge_paragraph,
        base / "warning" / "edu-dense-slide.pptx": build_edu_warn_dense_slide,
        base / "warning" / "edu-crowded-shurley.pdf": build_edu_warn_crowded_shurley,
        base / "failing" / "a4-page.pdf": build_fail_a4,
        base / "failing" / "unsafe-edge.pdf": build_fail_unsafe_edge,
        base / "failing" / "blank-final-page.pdf": build_fail_blank_final_page,
        base / "failing" / "mixed-sizes.pdf": build_fail_mixed_sizes,
        base / "failing" / "placeholder.pdf": build_fail_placeholder,
        base / "failing" / "html-no-print.html": build_fail_html_no_print,
        base / "failing" / "html-overflow.html": build_fail_html_overflow,
        base / "failing" / "docx-a4.docx": build_fail_docx_a4,
        base / "failing" / "pptx-out-of-bounds.pptx": build_fail_pptx_out_of_bounds,
        base / "failing" / "shurley-wrap.pdf": build_fail_shurley_wrap,
    }
    for target, builder in builders.items():
        target.parent.mkdir(parents=True, exist_ok=True)
        builder(target)
