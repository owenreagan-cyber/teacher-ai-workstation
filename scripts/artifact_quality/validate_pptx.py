from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import CheckStatus, PLACEHOLDER_PATTERNS, PreflightReport
from .profiles import ArtifactProfile


def _slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        chunks.append(shape.text)
    return "\n".join(chunks)


def validate_pptx(path: Path, profile: ArtifactProfile, report: PreflightReport) -> None:
    if not path.is_file():
        report.add(CheckStatus.FAIL, "PPTX input file does not exist")
        return

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        report.add(CheckStatus.FAIL, "PPTX file is not readable", details=str(exc))
        return

    report.add(CheckStatus.PASS, "PPTX opens successfully")
    slide_count = len(presentation.slides)
    if slide_count == 0:
        report.add(CheckStatus.FAIL, "Presentation has zero slides")
        return
    report.add(CheckStatus.PASS, f"{slide_count} slide{'s' if slide_count != 1 else ''} detected")

    slide_width = presentation.slide_width.inches
    slide_height = presentation.slide_height.inches
    if profile.artifact_type == "pptx-print":
        expected_w = profile.paper.width_points / 72.0
        expected_h = profile.paper.height_points / 72.0
        if abs(slide_width - expected_w) < 0.2 and abs(slide_height - expected_h) < 0.2:
            report.add(CheckStatus.PASS, "Slide dimensions match printable Letter profile")
        else:
            report.add(CheckStatus.FAIL, "Slide dimensions do not match printable profile", details=f"{slide_width:.2f}x{slide_height:.2f}in")
    elif profile.slides:
        if abs(slide_width - profile.slides.width_inches) < 0.3 and abs(slide_height - profile.slides.height_inches) < 0.3:
            report.add(CheckStatus.PASS, "Slide ratio matches projector profile")
        else:
            report.add(CheckStatus.WARN, "Slide ratio differs from projector profile", details=f"{slide_width:.2f}x{slide_height:.2f}in")

    min_title = profile.slides.min_title_font_pt if profile.slides else 28.0
    min_body = profile.slides.min_body_font_pt if profile.slides else 18.0
    overflow_slides: list[int] = []
    out_of_bounds: list[int] = []
    empty_slides: list[int] = []
    dense_slides: list[int] = []

    for index, slide in enumerate(presentation.slides, start=1):
        text = _slide_text(slide).strip()
        if not text:
            empty_slides.append(index)
        if len(text) > 900:
            dense_slides.append(index)

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            left = shape.left.inches if shape.left is not None else 0
            top = shape.top.inches if shape.top is not None else 0
            width = shape.width.inches if shape.width is not None else 0
            height = shape.height.inches if shape.height is not None else 0
            if left < 0 or top < 0 or left + width > slide_width + 0.05 or top + height > slide_height + 0.05:
                out_of_bounds.append(index)
            frame = shape.text_frame
            if frame is not None and getattr(frame, "auto_size", None) is None:
                if len(text) > 500 and height < 1.0:
                    overflow_slides.append(index)
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    size = run.font.size.pt if run.font.size is not None else None
                    if size is not None and size < min_body and paragraph.text.strip():
                        report.add(CheckStatus.WARN, f"Slide {index} body font below minimum", details=f"{size} pt")

    if empty_slides:
        report.add(CheckStatus.WARN, f"Empty slides detected: {empty_slides}")
    else:
        report.add(CheckStatus.PASS, "No empty slides detected")

    if out_of_bounds:
        report.add(CheckStatus.FAIL, f"Objects outside slide bounds on slides: {sorted(set(out_of_bounds))}")
    else:
        report.add(CheckStatus.PASS, "All objects appear within slide bounds")

    if overflow_slides:
        report.add(CheckStatus.WARN, f"Possible text overflow on slides: {sorted(set(overflow_slides))}")

    if dense_slides:
        report.add(CheckStatus.WARN, f"High text density on slides: {sorted(set(dense_slides))}", details="One-core-idea heuristic: review for projector readability.")

    full_text = "\n".join(_slide_text(slide) for slide in presentation.slides)
    placeholders = [pattern for pattern in PLACEHOLDER_PATTERNS if pattern.upper() in full_text.upper()]
    if placeholders:
        report.add(CheckStatus.FAIL, "Unresolved placeholders detected in PPTX", details=", ".join(placeholders))
    else:
        report.add(CheckStatus.PASS, "No unresolved placeholders detected in PPTX")

    if profile.artifact_type == "pptx-print":
        report.add(
            CheckStatus.WARN,
            "PPTX printable export not converted to PDF automatically",
            details="Export to PDF and run PDF preflight for final print proof.",
        )
